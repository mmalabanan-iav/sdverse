import os
import re
import json
import subprocess
from pathlib import Path
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

# -------------------------
# CONFIG
# -------------------------
OUTPUT_FOLDER = "output"
MODEL = "llama3"

QUESTIONS = [
    "provide short description of the product or services?",
    "provide long description of the product or services?"
]

# -------------------------
# HELPERS
# -------------------------

def check_ollama():
    """Check if Ollama is installed."""
    try:
        subprocess.run(["ollama", "--version"], check=True, capture_output=True)
        return True
    except Exception:
        print("Ollama is not installed or not in PATH.")
        return False

# ---- File parsers ----
def extract_text_from_pdf(file_path):
    text = []
    reader = PdfReader(file_path)
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)

def extract_text_from_pptx(file_path):
    text = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_docx(file_path):
    text = []
    doc = Document(file_path)
    for para in doc.paragraphs:
        text.append(para.text)
    return "\n".join(text)

def extract_text_from_file(file_path):
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

# ---- Web parser ----
def extract_text_from_url(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")

        # remove scripts and styles
        for tag in soup(["script", "style", "noscript"]):
            tag.extract()

        text_parts = []
        # title
        if soup.title:
            text_parts.append(soup.title.get_text())
        # headings and paragraphs
        for tag in soup.find_all(["h1","h2","h3","h4","h5","h6","p","li"]):
            text = tag.get_text(strip=True)
            if text:
                text_parts.append(text)

        return "\n".join(text_parts)
    except Exception as e:
        print(f"Failed to fetch {url}: {e}")
        return ""

def slugify_url(url):
    parsed = urlparse(url)
    slug = parsed.netloc.replace(".", "_") + parsed.path.replace("/", "_")
    slug = re.sub(r"[^a-zA-Z0-9_]+", "", slug)
    return slug.strip("_") or "link"

# ---- Ollama ----
def ask_ollama(question, context_text):
    """Send a question with context to Ollama and return the response."""
    prompt = f"Context:\n{context_text}\n\nQuestion: {question}"
    
    try:
        result = subprocess.run(
            ["ollama", "run", MODEL],
            input=prompt,  # Pass as string, not bytes
            capture_output=True,
            text=True,     # Text mode - expects string input
            timeout=300    # 5 minute timeout for long responses
        )
        
        if result.returncode != 0:
            print(f"Ollama error: {result.stderr}")
            return f"Error: {result.stderr}"
        
        return result.stdout.strip()
    
    except subprocess.TimeoutExpired:
        print(f"Ollama request timed out for question: {question}")
        return "Error: Request timed out"
    except Exception as e:
        print(f"Error calling Ollama: {e}")
        return f"Error: {str(e)}"

# ---- Processing ----
def process_file(file_path, output_folder):
    print(f"Processing file: {file_path.name}")
    
    try:
        text = extract_text_from_file(file_path)
        
        if not text.strip():
            print(f"No text extracted from {file_path.name}")
            return
        
        print(f"   Extracted {len(text)} characters")
        
        qa_results = {}
        for q in QUESTIONS:
            print(f"   Asking: {q}")
            answer = ask_ollama(q, text)
            qa_results[q] = answer

        output_data = {
            "source": str(file_path.name),
            "qa": qa_results
        }
        
        output_path = output_folder / f"{file_path.stem}.json"
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(output_data, f, indent=2, ensure_ascii=False)
        
        print(f"Saved results to {output_path}")
        
    except Exception as e:
        print(f"Error processing {file_path.name}: {e}")

def process_url(url, output_folder):
    print(f"Processing URL: {url}")
    
    try:
        text = extract_text_from_url(url)
        
        if not text.strip():
            print(f"No text extracted from {url}")
            return
        
        print(f"   Extracted {len(text)} characters")
        
        qa_results = {}
        for q in QUESTIONS:
            print(f"   Asking: {q}")
            answer = ask_ollama(q, text)
            qa_results[q] = answer

        output_data = {
            "source": url,
            "qa": qa_results
        }
        
        slug = slugify_url(url)
        output_path = output_folder / f"{slug}.json"
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(output_data, f, indent=2, ensure_ascii=False)
        
        print(f"Saved results to {output_path}")
        
    except Exception as e:
        print(f"Error processing {url}: {e}")

# -------------------------
# MAIN
# -------------------------
def main():
    if not check_ollama():
        return

    output_folder = Path(OUTPUT_FOLDER)
    output_folder.mkdir(exist_ok=True)

    while True:
        user_input = input("\nEnter a folder path (for files) or a link (http/https): ").strip()

        if user_input.startswith("http://") or user_input.startswith("https://"):
            process_url(user_input, output_folder)
        else:
            input_folder = Path(user_input)
            if not input_folder.exists() or not input_folder.is_dir():
                print(f"Input path {user_input} is not a valid folder.")
                continue
            
            files = [f for f in input_folder.iterdir() if f.suffix.lower() in [".pdf", ".pptx", ".docx"]]
            if not files:
                print("No supported files found in input folder.")
                continue
            
            print(f"Found {len(files)} file(s) to process")
            
            for file_path in files:
                process_file(file_path, output_folder)

        run_again = input("\n🔄 Do you want to process another folder/link? (y/n): ").strip().lower()
        if run_again != "y":
            print("Exiting. Goodbye!")
            break

if __name__ == "__main__":
    main()